"""
SARTHI — Round 1 PPT builder.

Generates SARTHI.pptx (16:9, 10 slides) from the locked slide copy in PPT-SLIDES.md.
Open the output in PowerPoint / Google Slides, or upload to Canva (Uploads → Import PPTX)
to polish visuals. Text/structure/brand colors ship pre-formatted.

Run: python build_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ---------------------------------------------------------------------------
# Brand palette
# ---------------------------------------------------------------------------
INDIGO_DEEP = RGBColor(0x1E, 0x1B, 0x4B)     # primary
INDIGO_MID  = RGBColor(0x3B, 0x37, 0x8A)
SAFFRON     = RGBColor(0xF2, 0x86, 0x2F)     # accent
OFFWHITE    = RGBColor(0xFA, 0xFA, 0xFB)
GRAY_TEXT   = RGBColor(0x4B, 0x4B, 0x57)
GRAY_LIGHT  = RGBColor(0xE8, 0xE8, 0xEF)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "Calibri"         # Inter/Manrope if installed, else Calibri
BODY_FONT  = "Calibri"


# ---------------------------------------------------------------------------
# Presentation setup (16:9)
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def add_bg(slide, color):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.shadow.inherit = False
    return rect


def add_rect(slide, left, top, width, height, fill=None, line=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.shadow.inherit = False
    if fill is None:
        rect.fill.background()
    else:
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(0.75)
    return rect


def add_rounded(slide, left, top, width, height, fill, line=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.shadow.inherit = False
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(0.75)
    return rect


def add_text(slide, left, top, width, height,
             text, *, size=18, bold=False, color=INDIGO_DEEP,
             font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, bullets, *,
                size=16, color=INDIGO_DEEP, font=BODY_FONT, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = "•  " + b
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def set_cell(cell, text, *, size=12, bold=False, color=INDIGO_DEEP,
             bg=None, align=PP_ALIGN.LEFT, font=BODY_FONT):
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    tf = cell.text_frame
    tf.clear()
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------------------
# SLIDE 1 — COVER
# ---------------------------------------------------------------------------
def slide_1_cover():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, INDIGO_DEEP)

    # Saffron accent bar on left
    add_rect(s, 0, Inches(3.2), Inches(0.25), Inches(1.1), fill=SAFFRON)

    # Big SARTHI wordmark
    add_text(s, Inches(0.8), Inches(2.8), Inches(11), Inches(1.5),
             "SARTHI", size=100, bold=True, color=WHITE, font=TITLE_FONT)

    # Tagline
    add_text(s, Inches(0.85), Inches(4.3), Inches(11), Inches(0.6),
             "Your AI Sarthi — from dream to degree.",
             size=26, color=SAFFRON, italic=True, font=BODY_FONT)

    # Round info
    add_text(s, Inches(0.85), Inches(5.6), Inches(11), Inches(0.4),
             "Round 1  ·  Problem Statement 2",
             size=16, color=WHITE, font=BODY_FONT)
    add_text(s, Inches(0.85), Inches(6.05), Inches(11), Inches(0.4),
             "CRP TenzorX 2026  ·  Poonawalla Fincorp National AI Hackathon",
             size=14, color=GRAY_LIGHT, font=BODY_FONT)

    # Team strip at bottom
    add_text(s, Inches(0.85), Inches(6.7), Inches(11), Inches(0.4),
             "Team: [Your Name]  ·  AI-augmented build",
             size=12, color=GRAY_LIGHT, italic=True, font=BODY_FONT)

    add_notes(s, "We're SARTHI — an agentic AI platform that guides "
                 "Tier-2/3 Indian students from the first dream of studying "
                 "abroad all the way to loan disbursement. In their language. "
                 "With memory. Zero human intervention.")


# ---------------------------------------------------------------------------
# Shared header for content slides
# ---------------------------------------------------------------------------
def content_header(slide, title):
    # Saffron tick
    add_rect(slide, Inches(0.5), Inches(0.55), Inches(0.12), Inches(0.55),
             fill=SAFFRON)
    add_text(slide, Inches(0.75), Inches(0.45), Inches(12), Inches(0.8),
             title, size=32, bold=True, color=INDIGO_DEEP, font=TITLE_FONT)
    # Thin rule under title
    add_rect(slide, Inches(0.5), Inches(1.25), Inches(12.3), Emu(12700),
             fill=GRAY_LIGHT)


def footer_strip(slide, text="SARTHI  ·  Your AI Sarthi — from dream to degree."):
    add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), fill=INDIGO_DEEP)
    add_text(slide, Inches(0.5), Inches(7.18), Inches(12.3), Inches(0.3),
             text, size=10, color=WHITE, italic=True, font=BODY_FONT,
             align=PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# SLIDE 2 — THE PROBLEM
# ---------------------------------------------------------------------------
def slide_2_problem():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, OFFWHITE)
    content_header(s, "Meet Priya. India has 10 million Priyas.")

    # LEFT panel — Priya profile card
    add_rounded(s, Inches(0.5), Inches(1.6), Inches(5.3), Inches(5.2),
                fill=INDIGO_DEEP)
    add_text(s, Inches(0.85), Inches(1.85), Inches(4.6), Inches(0.5),
             "PRIYA SHARMA", size=14, bold=True, color=SAFFRON, font=TITLE_FONT)
    add_text(s, Inches(0.85), Inches(2.35), Inches(4.6), Inches(0.5),
             "21 · Nagpur", size=20, bold=True, color=WHITE, font=TITLE_FONT)

    profile_bullets = [
        "Final-year Mech Engg, VNIT Nagpur · CGPA 7.8",
        "Father: small auto-parts business · ₹6L / yr",
        "Dream: MS Robotics, US / Canada",
        "Budget: ₹40–60L · ~70% on loan",
        "Speaks Marathi–Hindi at home, English in class",
        "First-in-family to consider studying abroad",
    ]
    tb = s.shapes.add_textbox(Inches(0.85), Inches(3.05), Inches(4.6), Inches(3.7))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(profile_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.3
        run = p.add_run()
        run.text = "— " + b
        run.font.name = BODY_FONT
        run.font.size = Pt(13)
        run.font.color.rgb = WHITE

    # RIGHT panel — 7 pains grid (2 cols × 4 rows, one cell empty)
    pains = [
        ("🤯", "Which test? IELTS, GRE, TOEFL?"),
        ("🧭", "500 US unis — which 10 fit me?"),
        ("😰", "Will I get in? Worth ₹50L?"),
        ("🏠", "Family can't guide me."),
        ("🗣️", "Think in Marathi. Internet in English."),
        ("🏦", "Messy papers — any NBFC say yes?"),
        ("💸", "Counselors ₹1L — push their favorites."),
    ]
    add_text(s, Inches(6.1), Inches(1.6), Inches(6.8), Inches(0.5),
             "Her 7 real pains",
             size=16, bold=True, color=INDIGO_DEEP, font=TITLE_FONT)

    col_w, row_h = Inches(3.3), Inches(1.2)
    start_x, start_y = Inches(6.1), Inches(2.15)
    for idx, (emoji, text) in enumerate(pains):
        col = idx % 2
        row = idx // 2
        x = start_x + col_w * col
        y = start_y + row_h * row
        add_rounded(s, x, y, Inches(3.2), Inches(1.1),
                    fill=WHITE, line=GRAY_LIGHT)
        add_text(s, x + Inches(0.15), y + Inches(0.15),
                 Inches(0.5), Inches(0.8),
                 emoji, size=22, color=INDIGO_DEEP, font=BODY_FONT)
        add_text(s, x + Inches(0.75), y + Inches(0.2),
                 Inches(2.3), Inches(0.9),
                 text, size=11, color=GRAY_TEXT, font=BODY_FONT,
                 anchor=MSO_ANCHOR.TOP)

    footer_strip(s)
    add_notes(s, "Priya isn't one student — she represents 10 million "
                 "Indian students every year who face a broken study-abroad "
                 "journey. These are the 7 pains no platform solves together today.")


# ---------------------------------------------------------------------------
# SLIDE 3 — OUR INSIGHT
# ---------------------------------------------------------------------------
def slide_3_insight():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, OFFWHITE)
    content_header(s, "Relationship  >  Transaction.")

    # Today's NBFC — gray row
    add_rounded(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.9),
                fill=GRAY_LIGHT)
    add_text(s, Inches(0.9), Inches(1.9), Inches(2.5), Inches(0.5),
             "TODAY'S NBFC", size=14, bold=True, color=GRAY_TEXT, font=TITLE_FONT)
    add_text(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.1),
             "Shows up on Day-330. Fights a commodity price war on loan rates.\n"
             "Student has no loyalty. NBFC has no engagement. Loan is a transaction.",
             size=15, color=GRAY_TEXT, font=BODY_FONT)

    # SARTHI — indigo row
    add_rounded(s, Inches(0.5), Inches(3.85), Inches(12.3), Inches(1.9),
                fill=INDIGO_DEEP)
    add_text(s, Inches(0.9), Inches(4.05), Inches(2.5), Inches(0.5),
             "SARTHI", size=14, bold=True, color=SAFFRON, font=TITLE_FONT)
    add_text(s, Inches(0.9), Inches(4.55), Inches(11.5), Inches(1.1),
             "Shows up on Day-1. Becomes the 12-month AI mentor.\n"
             "The loan is the natural outcome — not the start.",
             size=15, color=WHITE, font=BODY_FONT)

    # Reframe line
    add_text(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.7),
             "We don't sell loans. We build 12-month relationships that end in loans.",
             size=22, bold=True, color=SAFFRON, italic=True, font=TITLE_FONT,
             align=PP_ALIGN.CENTER)

    footer_strip(s)
    add_notes(s, "Every NBFC is fighting for the last 30 days of a student's "
                 "12-month journey. SARTHI shows up on Day 1. We become the mentor. "
                 "The loan becomes the natural outcome, not a price war.")


# ---------------------------------------------------------------------------
# SLIDE 4 — MEET SARTHI
# ---------------------------------------------------------------------------
def slide_4_meet_sarthi():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, OFFWHITE)
    content_header(s, "Not a chatbot. An agent.")

    # Product identity strip
    add_text(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.5),
             "SARTHI (सारथी) — the charioteer  ·  Your AI Sarthi, from dream to degree.",
             size=15, bold=True, color=INDIGO_DEEP, italic=False, font=TITLE_FONT)
    add_text(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.5),
             "Agentic AI platform guiding Tier-2/3 Indian students end-to-end — "
             "in their language, with memory, with action.",
             size=12, color=GRAY_TEXT, italic=True, font=BODY_FONT)

    # 3×3 competitive matrix
    rows = [
        ("Category", "Who they are", "Their gap", "SARTHI's answer"),
        ("Counselors", "Leverage Edu · Yocket · Shiksha",
         "₹50k–₹2L · biased · English-only",
         "Free · consistent AI · vernacular · unbiased"),
        ("Loan NBFCs", "Credila · Avanse · Auxilo · Prodigy",
         "Arrive too late · commodity price war",
         "12-month pre-loan relationship · propensity data"),
        ("Generic AI", "ChatGPT · Gemini · other LLMs",
         "No memory · no action · no India · no loan",
         "Memory · tools · Bharat · loan-native"),
    ]
    tbl_left, tbl_top = Inches(0.5), Inches(2.7)
    tbl_w, tbl_h = Inches(12.3), Inches(3.4)
    table_shape = s.shapes.add_table(len(rows), 4, tbl_left, tbl_top, tbl_w, tbl_h)
    table = table_shape.table
    table.columns[0].width = Inches(1.9)
    table.columns[1].width = Inches(3.1)
    table.columns[2].width = Inches(3.5)
    table.columns[3].width = Inches(3.8)

    for r, row in enumerate(rows):
        is_header = (r == 0)
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            if is_header:
                set_cell(cell, val, size=12, bold=True, color=WHITE,
                         bg=INDIGO_DEEP, align=PP_ALIGN.LEFT)
            else:
                bg = WHITE if (r % 2 == 1) else GRAY_LIGHT
                is_first = (c == 0)
                set_cell(cell, val, size=11,
                         bold=is_first, color=INDIGO_DEEP if is_first else GRAY_TEXT,
                         bg=bg, align=PP_ALIGN.LEFT)

    # Tag line
    add_text(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.6),
             "Counselors charge too much. Loan NBFCs show up too late. "
             "Generic AI doesn't know India. SARTHI solves all three.",
             size=15, bold=True, color=SAFFRON, italic=True,
             align=PP_ALIGN.CENTER, font=TITLE_FONT)

    footer_strip(s)
    add_notes(s, "SARTHI sits at the intersection of three broken categories. "
                 "Counselors overcharge and bias students. Loan NBFCs arrive too "
                 "late and commoditize. Generic AI has no memory, no India context, "
                 "no loan funnel. SARTHI is the only agent that covers all three gaps.")


# ---------------------------------------------------------------------------
# SLIDE 5 — HERO: THE 4-PHASE JOURNEY
# ---------------------------------------------------------------------------
def slide_5_journey():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, OFFWHITE)
    content_header(s, "Priya's journey. 7 features. One agent.")

    phases = [
        ("01", "DISCOVER",
         "IG reel → voice chat in Hinglish. Agent builds profile + "
         "personalized application timeline (IELTS → GRE → Apps → Visa).",
         "F1"),
        ("02", "DECIDE",
         "Shortlist 10 universities · admission probability · "
         "per-uni ROI (cost vs salary vs EMI).",
         "F2 + F3"),
        ("03", "APPLY",
         "Socratic SOP Co-Pilot references her actual Bajaj internship — "
         "because it remembers.",
         "F4"),
        ("04", "FUND",
         "\"You qualify for ₹45L at 10.2%.\" Auto-fills 90% of the loan "
         "app from past conversations.",
         "F5 + F6"),
        ("↻", "AMPLIFY",
         "Shareable \"Study Abroad Passport\" → 3 friends join → "
         "loop closes.",
         "F7"),
    ]

    card_w = Inches(2.45)
    gap    = Inches(0.08)
    start_x = Inches(0.5)
    card_top = Inches(1.7)
    card_h = Inches(4.4)

    for i, (num, name, text, fids) in enumerate(phases):
        x = start_x + (card_w + gap) * i
        # Card body
        bg_color = INDIGO_DEEP if i == len(phases) - 1 else WHITE
        line_color = INDIGO_DEEP if i == len(phases) - 1 else GRAY_LIGHT
        add_rounded(s, x, card_top, card_w, card_h, fill=bg_color, line=line_color)

        # Phase number (big)
        add_text(s, x, card_top + Inches(0.2), card_w, Inches(0.7),
                 num, size=36, bold=True,
                 color=SAFFRON if i == len(phases) - 1 else SAFFRON,
                 font=TITLE_FONT, align=PP_ALIGN.CENTER)
        # Phase name
        add_text(s, x, card_top + Inches(1.0), card_w, Inches(0.5),
                 name, size=14, bold=True,
                 color=WHITE if i == len(phases) - 1 else INDIGO_DEEP,
                 font=TITLE_FONT, align=PP_ALIGN.CENTER)
        # Body text
        add_text(s, x + Inches(0.2), card_top + Inches(1.6), card_w - Inches(0.4), Inches(2.2),
                 text, size=11,
                 color=GRAY_LIGHT if i == len(phases) - 1 else GRAY_TEXT,
                 font=BODY_FONT, align=PP_ALIGN.LEFT)
        # Feature IDs badge
        add_rounded(s, x + Inches(0.25), card_top + card_h - Inches(0.7),
                    card_w - Inches(0.5), Inches(0.45),
                    fill=SAFFRON)
        add_text(s, x + Inches(0.25), card_top + card_h - Inches(0.67),
                 card_w - Inches(0.5), Inches(0.4),
                 fids, size=12, bold=True, color=INDIGO_DEEP,
                 font=TITLE_FONT, align=PP_ALIGN.CENTER)

    # Hero banner line
    add_rounded(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.65),
                fill=INDIGO_DEEP)
    add_text(s, Inches(0.5), Inches(6.33), Inches(12.3), Inches(0.6),
             "4 phases · One agent · Zero human intervention · One ₹45L loan disbursed.",
             size=17, bold=True, color=SAFFRON, italic=True,
             font=TITLE_FONT, align=PP_ALIGN.CENTER)

    footer_strip(s)
    add_notes(s, "This is the heart of the pitch. Four phases — Discover, Decide, "
                 "Apply, Fund — each one powered by AI features, each one seamlessly "
                 "handing off to the next. The loop closes when Priya shares her "
                 "Study Abroad Passport and three friends join. This is what zero "
                 "human intervention actually looks like.")


# ---------------------------------------------------------------------------
# SLIDE 6 — AI ARCHITECTURE
# ---------------------------------------------------------------------------
def slide_6_architecture():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, OFFWHITE)
    content_header(s, "AI architecture — agent, not chatbot.")

    # Main layered diagram on the left (8 inches wide)
    x0 = Inches(0.5)
    layer_w = Inches(8.4)
    layer_h = Inches(0.85)
    y = Inches(1.55)

    layers = [
        ("Frontend", "Next.js · shadcn/ui · Tailwind · PWA", INDIGO_MID, WHITE),
        ("Agent Orchestrator",
         "LLM-agnostic · long-term memory · tool-calling · multi-model fallback",
         INDIGO_DEEP, WHITE),
    ]
    for label, desc, bg, fg in layers:
        add_rounded(s, x0, y, layer_w, layer_h, fill=bg)
        add_text(s, x0 + Inches(0.25), y + Inches(0.1), layer_w, Inches(0.4),
                 label, size=14, bold=True, color=SAFFRON, font=TITLE_FONT)
        add_text(s, x0 + Inches(0.25), y + Inches(0.42), layer_w - Inches(0.3), Inches(0.4),
                 desc, size=11, color=fg, font=BODY_FONT)
        y += layer_h + Inches(0.08)

    # Tools row: 5 small boxes
    tools = [
        "University\nShortlister",
        "ROI\nPredictor",
        "SOP\nCo-Pilot",
        "Loan Offer\n+ Eligibility",
        "Passport\nGenerator",
    ]
    tool_w = (layer_w - Inches(0.32)) / 5
    tool_h = Inches(0.95)
    for i, t in enumerate(tools):
        tx = x0 + (tool_w + Inches(0.08)) * i
        add_rounded(s, tx, y, tool_w, tool_h, fill=SAFFRON)
        add_text(s, tx, y + Inches(0.12), tool_w, tool_h,
                 t, size=11, bold=True, color=INDIGO_DEEP,
                 font=TITLE_FONT, align=PP_ALIGN.CENTER)
    y += tool_h + Inches(0.08)

    # Data layer
    add_rounded(s, x0, y, layer_w, layer_h, fill=INDIGO_MID)
    add_text(s, x0 + Inches(0.25), y + Inches(0.1), layer_w, Inches(0.4),
             "Data & Integration Layer", size=14, bold=True, color=SAFFRON,
             font=TITLE_FONT)
    add_text(s, x0 + Inches(0.25), y + Inches(0.42), layer_w - Inches(0.3), Inches(0.4),
             "Vector DB · Postgres · Whisper + Bhashini · Poonawalla Loan API",
             size=11, color=WHITE, font=BODY_FONT)

    # RIGHT column — chatbot vs agent comparison
    right_x = Inches(9.2)
    add_text(s, right_x, Inches(1.55), Inches(3.8), Inches(0.4),
             "Chatbot  vs.  SARTHI Agent",
             size=13, bold=True, color=INDIGO_DEEP, font=TITLE_FONT)

    compare = [
        ("Stateless",        "Memory per user"),
        ("Answers",          "Takes actions"),
        ("English only",     "Vernacular voice"),
        ("No tools",         "Tool-calling"),
    ]
    cy = Inches(2.05)
    for l, r in compare:
        add_rounded(s, right_x, cy, Inches(1.85), Inches(0.55),
                    fill=GRAY_LIGHT)
        add_text(s, right_x + Inches(0.1), cy + Inches(0.08),
                 Inches(1.65), Inches(0.45),
                 l, size=11, color=GRAY_TEXT, font=BODY_FONT,
                 align=PP_ALIGN.CENTER)
        add_rounded(s, right_x + Inches(1.95), cy, Inches(1.85), Inches(0.55),
                    fill=INDIGO_DEEP)
        add_text(s, right_x + Inches(2.05), cy + Inches(0.08),
                 Inches(1.65), Inches(0.45),
                 r, size=11, bold=True, color=SAFFRON, font=BODY_FONT,
                 align=PP_ALIGN.CENTER)
        cy += Inches(0.65)

    # Compliance strip (saffron bar)
    add_rounded(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.65),
                fill=SAFFRON)
    add_text(s, Inches(0.5), Inches(6.33), Inches(12.3), Inches(0.6),
             "Consent-first  ·  DPDP-compliant  ·  RBI digital-lending aligned  ·  Explainable by default",
             size=14, bold=True, color=INDIGO_DEEP, italic=True,
             font=TITLE_FONT, align=PP_ALIGN.CENTER)

    footer_strip(s)
    add_notes(s, "We built SARTHI as a true agent. It remembers every conversation. "
                 "It calls specific tools — ML shortlister, ROI calculator, SOP "
                 "co-pilot, loan offer engine. It's LLM-agnostic so we can optimize "
                 "for cost and latency. And we've engineered DPDP and RBI compliance "
                 "from day one — not as an afterthought.")


# ---------------------------------------------------------------------------
# SLIDE 7 — GROWTH LOOP
# ---------------------------------------------------------------------------
def slide_7_loop():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, OFFWHITE)
    content_header(s, "The zero-human-intervention growth loop.")

    # Subtitle / bonus callout
    add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.4),
             "Directly answering the hackathon's bonus challenge.",
             size=13, italic=True, color=GRAY_TEXT, font=BODY_FONT)

    # 5-node loop: laid out as pentagon / big arrows
    nodes = [
        ("ACQUIRE",  "AI-generated IG reels + SEO blogs\n(1 / day / city × course)"),
        ("ENGAGE",   "Free shortlister / ROI hook\nProfile built in 5 minutes"),
        ("NURTURE",  "Proactive WhatsApp nudges\nStreaks · milestones"),
        ("CONVERT",  "Personalized loan offer\nAuto-filled application"),
        ("AMPLIFY",  "Study Abroad Passport\n3 friends join → loop closes"),
    ]

    # Arrange as 5 cards in a row, with arrow from AMPLIFY to ACQUIRE visualized as text
    node_w = Inches(2.35)
    node_gap = Inches(0.08)
    start_x = Inches(0.55)
    top = Inches(2.2)
    h = Inches(2.6)

    for i, (title, body) in enumerate(nodes):
        x = start_x + (node_w + node_gap) * i
        add_rounded(s, x, top, node_w, h, fill=INDIGO_DEEP)
        # Number circle
        add_rounded(s, x + (node_w - Inches(0.7)) / 2, top + Inches(0.18),
                    Inches(0.7), Inches(0.7),
                    fill=SAFFRON)
        add_text(s, x + (node_w - Inches(0.7)) / 2, top + Inches(0.22),
                 Inches(0.7), Inches(0.65),
                 str(i + 1), size=22, bold=True, color=INDIGO_DEEP,
                 font=TITLE_FONT, align=PP_ALIGN.CENTER)
        add_text(s, x, top + Inches(1.0), node_w, Inches(0.5),
                 title, size=14, bold=True, color=SAFFRON,
                 font=TITLE_FONT, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), top + Inches(1.5), node_w - Inches(0.3), Inches(1.1),
                 body, size=10, color=WHITE, font=BODY_FONT,
                 align=PP_ALIGN.CENTER)

    # Loop back arrow — simple text callout
    add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.5),
             "↻  Amplify → Acquire  (closed loop)",
             size=13, bold=True, italic=True, color=SAFFRON,
             font=TITLE_FONT, align=PP_ALIGN.CENTER)

    # K-factor callout
    add_rounded(s, Inches(0.5), Inches(5.65), Inches(12.3), Inches(1.3),
                fill=INDIGO_DEEP)
    add_text(s, Inches(0.7), Inches(5.75), Inches(12), Inches(0.6),
             "Target K-factor: 1.5",
             size=22, bold=True, color=SAFFRON, font=TITLE_FONT,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.6),
             "Every 2 converted students bring 3 friends into the top of the funnel. "
             "Acquisition cost approaches zero over time.",
             size=13, color=WHITE, italic=True, font=BODY_FONT,
             align=PP_ALIGN.CENTER)

    footer_strip(s)
    add_notes(s, "This is exactly the bonus challenge — an AI growth loop where "
                 "acquisition, engagement, nurturing, conversion, and amplification "
                 "all happen with zero human intervention. Every Priya who gets a "
                 "loan brings three friends into the funnel. Our acquisition cost "
                 "approaches zero over time.")


# ---------------------------------------------------------------------------
# SLIDE 8 — BUSINESS IMPACT
# ---------------------------------------------------------------------------
def slide_8_business():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, OFFWHITE)
    content_header(s, "From ₹5,000 CAC to a 30-year relationship.")

    # Killer quote
    add_rounded(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(1.1),
                fill=INDIGO_DEEP)
    add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.0),
             "SARTHI converts a ₹5,000 CAC into a ₹500 CAC,\n"
             "and a one-time loan into a 30-year relationship.",
             size=20, bold=True, italic=True, color=SAFFRON,
             font=TITLE_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Revenue streams header
    add_text(s, Inches(0.5), Inches(2.75), Inches(12.3), Inches(0.4),
             "Four revenue streams — not just loan interest",
             size=14, bold=True, color=INDIGO_DEEP, font=TITLE_FONT)

    # 2x2 revenue tiles
    streams = [
        ("1.  Loan Origination (Primary)",
         "Revenue-share on every Poonawalla-disbursed loan.",
         "~1–2% of loan value",
         SAFFRON),
        ("2.  University Lead-Gen",
         "Universities pay per qualified, intent-verified applicant.",
         "₹500–₹2,000 per lead",
         INDIGO_MID),
        ("3.  Premium Student Tier",
         "Unlimited SOPs · priority nudges · mock-interview coach.",
         "₹299–₹499 / month",
         INDIGO_MID),
        ("4.  White-label Agent (Phase 2)",
         "License agent to other NBFCs / BFSI after exclusivity.",
         "Annual contract",
         SAFFRON),
    ]
    tile_w = Inches(6.0)
    tile_h = Inches(1.55)
    tile_gap = Inches(0.25)
    start_x = Inches(0.5)
    start_y = Inches(3.25)
    for i, (name, desc, money, accent) in enumerate(streams):
        col = i % 2
        row = i // 2
        x = start_x + (tile_w + tile_gap) * col
        y = start_y + (tile_h + tile_gap) * row
        add_rounded(s, x, y, tile_w, tile_h, fill=WHITE, line=GRAY_LIGHT)
        # accent bar
        add_rect(s, x, y, Inches(0.12), tile_h, fill=accent)
        add_text(s, x + Inches(0.35), y + Inches(0.12), tile_w - Inches(0.5), Inches(0.4),
                 name, size=13, bold=True, color=INDIGO_DEEP, font=TITLE_FONT)
        add_text(s, x + Inches(0.35), y + Inches(0.5), tile_w - Inches(0.5), Inches(0.55),
                 desc, size=11, color=GRAY_TEXT, font=BODY_FONT)
        add_text(s, x + Inches(0.35), y + Inches(1.05), tile_w - Inches(0.5), Inches(0.4),
                 money, size=12, bold=True, italic=True, color=SAFFRON, font=TITLE_FONT)

    # GTM strip
    add_rounded(s, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.5),
                fill=INDIGO_DEEP)
    add_text(s, Inches(0.5), Inches(6.48), Inches(12.3), Inches(0.45),
             "GTM: Launch exclusively with Poonawalla Fincorp  ·  "
             "24-month exclusivity  ·  Multi-lender expansion in Phase 2",
             size=12, bold=True, color=WHITE, italic=True,
             font=TITLE_FONT, align=PP_ALIGN.CENTER)

    footer_strip(s)
    add_notes(s, "SARTHI isn't a loan origination tool — it's a platform with four "
                 "revenue streams and a 30-year customer LTV. Loan rev-share is "
                 "primary, but lead-gen fees, premium subscriptions, and Phase 2 "
                 "white-label multiply our upside. Poonawalla is our launch partner, "
                 "not our product.")


# ---------------------------------------------------------------------------
# SLIDE 9 — ROADMAP & PROTOTYPE
# ---------------------------------------------------------------------------
def slide_9_roadmap():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, OFFWHITE)
    content_header(s, "Prototype May 3. Scale-ready by month 12.")

    columns = [
        ("14-DAY BUILD",
         "Apr 19 → May 3",
         [
             "W1: Agent core + memory",
             "W1: Shortlister + ROI Predictor",
             "W2: SOP Co-Pilot + Loan Offer",
             "W2: Auto-fill + Passport + Voice",
             "Deploy: Vercel + GitHub",
             "Scripted Priya demo flow",
         ],
         INDIGO_MID),
        ("PHASE 2",
         "6 months",
         [
             "Domestic PG market",
             "CAT / GATE · IIMs · NITs",
             "Private MBAs",
             "+5 Indian languages",
             "(Tamil · Telugu · Marathi · Bengali · Kannada)",
             "Real university data partnerships",
         ],
         INDIGO_DEEP),
        ("PHASE 3",
         "12 months",
         [
             "White-label agent for BFSI",
             "Licensed to other NBFCs",
             "Intl corridors: Nepal, BD, SL → abroad",
             "Cross-sell pipeline",
             "(personal · home · business loans)",
             "Data flywheel maturity",
         ],
         SAFFRON),
    ]

    col_w = Inches(4.0)
    gap = Inches(0.15)
    start_x = Inches(0.55)
    top = Inches(1.55)
    col_h = Inches(5.3)

    for i, (title, when, bullets, accent) in enumerate(columns):
        x = start_x + (col_w + gap) * i
        # Column background
        add_rounded(s, x, top, col_w, col_h, fill=WHITE, line=GRAY_LIGHT)
        # accent header bar
        add_rounded(s, x, top, col_w, Inches(0.9), fill=accent)
        add_text(s, x + Inches(0.3), top + Inches(0.1), col_w - Inches(0.4), Inches(0.4),
                 title, size=16, bold=True,
                 color=INDIGO_DEEP if accent == SAFFRON else WHITE,
                 font=TITLE_FONT)
        add_text(s, x + Inches(0.3), top + Inches(0.5), col_w - Inches(0.4), Inches(0.4),
                 when, size=12, italic=True,
                 color=INDIGO_DEEP if accent == SAFFRON else GRAY_LIGHT,
                 font=BODY_FONT)

        # Bullets
        by = top + Inches(1.1)
        for b in bullets:
            add_text(s, x + Inches(0.3), by, col_w - Inches(0.4), Inches(0.45),
                     "•  " + b, size=12, color=INDIGO_DEEP, font=BODY_FONT)
            by += Inches(0.55)

    # Bottom callout
    add_text(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
             "Commitment: working prototype + GitHub by May 3.",
             size=12, bold=True, italic=True, color=SAFFRON,
             font=TITLE_FONT, align=PP_ALIGN.CENTER)

    footer_strip(s)
    add_notes(s, "A working prototype with scripted Priya flow lands on GitHub by "
                 "May 3. Phase 1 launches with study-abroad only. By month 12, "
                 "we're white-labeling the agent to the entire Indian BFSI sector "
                 "and expanding to South Asian corridors.")


# ---------------------------------------------------------------------------
# SLIDE 10 — CLOSING
# ---------------------------------------------------------------------------
def slide_10_closing():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, INDIGO_DEEP)

    # Saffron accent bar
    add_rect(s, 0, Inches(1.8), Inches(0.25), Inches(1.0), fill=SAFFRON)

    add_text(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.5),
             "Every student deserves a Sarthi.",
             size=54, bold=True, color=WHITE,
             font=TITLE_FONT, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.6),
             "Your AI Sarthi — from dream to degree.",
             size=22, italic=True, color=SAFFRON,
             font=BODY_FONT, align=PP_ALIGN.CENTER)

    # Team block
    add_text(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(0.5),
             "Team: [Your Name]",
             size=16, bold=True, color=WHITE,
             font=TITLE_FONT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.05), Inches(12.3), Inches(0.5),
             "Human product owner  +  Claude-powered AI co-developer",
             size=12, italic=True, color=GRAY_LIGHT,
             font=BODY_FONT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.5),
             "Our build process is itself a proof of 'zero-human-intervention' feasibility.",
             size=11, italic=True, color=GRAY_LIGHT,
             font=BODY_FONT, align=PP_ALIGN.CENTER)

    # Contact
    add_text(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.5),
             "[your-email]   ·   github.com/[your-handle]   ·   linkedin.com/in/[you]",
             size=13, color=WHITE,
             font=BODY_FONT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.4),
             "Live prototype: coming May 3 (Round 2)",
             size=11, italic=True, color=SAFFRON,
             font=BODY_FONT, align=PP_ALIGN.CENTER)

    add_notes(s, "SARTHI isn't just another hackathon submission. It's the mentor "
                 "10 million Indian students have been waiting for. Thank you.")


# ---------------------------------------------------------------------------
# Build everything
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    slide_1_cover()
    slide_2_problem()
    slide_3_insight()
    slide_4_meet_sarthi()
    slide_5_journey()
    slide_6_architecture()
    slide_7_loop()
    slide_8_business()
    slide_9_roadmap()
    slide_10_closing()

    out_path = "SARTHI.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")
